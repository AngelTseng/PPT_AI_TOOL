from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

BASE_DIR = Path(__file__).resolve().parent
TEMPLATE = BASE_DIR / "template.pptx"
OUT = BASE_DIR / "template_map.json"

if not TEMPLATE.exists():
    raise FileNotFoundError(f"Template not found: {TEMPLATE}")

prs = Presentation(str(TEMPLATE))
data = []


def _normalize_shape_names(shape_names: set[str]) -> set[str]:
    return {str(n).strip().lower() for n in shape_names if str(n).strip()}

def detect_template_type(shape_names: set, shapes: list, slide_index: int, total_slides: int) -> str:
    names = _normalize_shape_names(shape_names)

    def has(name: str) -> bool:
        return name.strip().lower() in names

    def has_img(prefix: str = "img") -> bool:
        prefix = prefix.lower()
        return any(n == prefix or n.startswith(f"{prefix}_") for n in names)

    if has("topic") and has("speaker_name"):
        return "cover"

    if has("outline") and has("agenda_1"):
        return "agenda"

    if has("agenda_name"):
        return "section"

    if has("flow_chart_1") or has("flow_chart_2") or has("flow_chart_3"):
        return "flow"

    if slide_index == total_slides:
        return "end"

    for shp in shapes:
        if str(shp.get("name", "")).strip().lower() == "sheet_1" and shp.get("has_table"):
            return "table"

    if has("title") and has("content") and has_img("img"):
        return "content_image"

    if has("title") and has("content") and not has_img("img"):
        return "content_text"

    if has("content_1") and has("content_2") and has("content_3") and has("content_4"):
        if has("item_1") and has("item_2") and has("item_3") and has("item_4"):
            return "content_4_b"
        return "content_4_a"

    if (
        has("item_1") and has("item_2") and has("item_3")
        and has("content_1") and has("content_2") and has("content_3")
    ):
        # 先判特徵版型，再判通用型
        if has_img("img"):
            return "content_3extra_image"

        if any(str(s.get("name", "")).strip().lower().startswith("straight connector") for s in shapes):
            return "content_3extra_a"

        rounded_rects = sum(
            1 for s in shapes
            if "矩形: 圓角" in str(s.get("name", "")).strip()
        )
        if rounded_rects >= 3:
            return "content_3extra_b"

        return "content_3extra"

    if (
        has("item_1") and has("item_2")
        and has("content_1") and has("content_2")
        and not has("item_3") and not has("content_3")
    ):
        return "content_2_c"

    if (
        has("title")
        and has("content_1") and has("content_2")
        and has("img_1") and has("img_2")
        and (has("title_content_1") or has("item_1"))
        and (has("title_content_2") or has("item_2"))
    ):
        return "content_2_a"

    if (
        has("content_1") and has("content_2")
        and has("img_1") and has("img_2")
        and (has("title_content_1") or has("item_1"))
        and (has("title_content_2") or has("item_2"))
    ):
        return "content_2_b"

    if (
        has("content_1") and has("content_2")
        and has("title_content_1") and has("title_content_2")
        and not has("img_1") and not has("img_2")
    ):
        return "content_2_text"

    return "unknown"


def classify_shape(shp, slide_type: str, slide_index: int):
    name = str(getattr(shp, "name", "") or "").strip()
    lname = name.lower()

    has_text = bool(getattr(shp, "has_text_frame", False))
    has_table = bool(getattr(shp, "has_table", False))
    has_chart = bool(getattr(shp, "has_chart", False))
    shape_type = getattr(shp, "shape_type", None)

    protected_text_names = {
        "topic",
        "speaker_name",
        "outline",
        "agenda_name",
        "agenda_1", "agenda_2", "agenda_3", "agenda_4", "agenda_5",
        "title",
        "content",
        "content_1", "content_2", "content_3", "content_4",
        "item_1", "item_2", "item_3", "item_4",
        "title_content_1", "title_content_2",
    }

    if has_table:
        return "protected", False, "table"
    if has_chart:
        return "protected", False, "chart"

    if lname.startswith("flow_chart_"):
        return "protected", False, "flow_chart"

    if lname in protected_text_names:
        return "protected", False, "text_container"

    if slide_type == "cover" and shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "protected", False, "cover_image"

    if (
        slide_index in {4, 10}
        and slide_type in {"content_image", "content_3extra_image"}
        and shape_type == MSO_SHAPE_TYPE.PICTURE
    ):
        if lname in {"img", "img_1", "main_image", "picture_1"} or lname.startswith("img"):
            return "image", False, "replaceable_main_picture"

    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "background", True, "decorative_picture"

    if shape_type in {
        MSO_SHAPE_TYPE.AUTO_SHAPE,
        MSO_SHAPE_TYPE.FREEFORM,
        MSO_SHAPE_TYPE.LINE,
        MSO_SHAPE_TYPE.GROUP,
        MSO_SHAPE_TYPE.TEXT_EFFECT,
    }:
        return "background", True, "decorative_shape"

    if has_text:
        return "background", True, "decorative_shape_with_text_frame"

    return "protected", False, "unknown_object"

for i, slide in enumerate(prs.slides, start=1):
    raw_shapes = []
    for shp in slide.shapes:
        raw_shapes.append({
            "name": str(getattr(shp, "name", "") or "").strip(),
            "is_placeholder": getattr(shp, "is_placeholder", False),
            "has_text": getattr(shp, "has_text_frame", False),
            "has_table": getattr(shp, "has_table", False),
            "has_chart": getattr(shp, "has_chart", False),
            "shape_type": str(getattr(shp, "shape_type", "")),
        })

    shape_names = {s["name"] for s in raw_shapes}
    detected_type = detect_template_type(
        shape_names=shape_names,
        shapes=raw_shapes,
        slide_index=i,
        total_slides=len(prs.slides),
    )

    enriched_shapes = []
    for shp in slide.shapes:
        role, allow_text_overlap, protect_reason = classify_shape(
            shp=shp,
            slide_type=detected_type,
            slide_index=i,
        )
        enriched_shapes.append({
            "name": str(getattr(shp, "name", "") or "").strip(),
            "is_placeholder": getattr(shp, "is_placeholder", False),
            "has_text": getattr(shp, "has_text_frame", False),
            "has_table": getattr(shp, "has_table", False),
            "has_chart": getattr(shp, "has_chart", False),
            "shape_type": str(getattr(shp, "shape_type", "")),
            "role": role,
            "allow_text_overlap": allow_text_overlap,
            "protect_reason": protect_reason,
        })

    data.append({
        "slide_index": i,
        "layout": slide.slide_layout.name,
        "detected_type": detected_type,
        "shapes": enriched_shapes,
    })

with open(OUT, "w", encoding="utf-8") as f:
    json.dump(data, f, ensure_ascii=False, indent=2)

print(f"Saved: {OUT}")